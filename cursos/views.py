from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth.forms import UserCreationForm, AuthenticationForm
from django.contrib.auth.models import User
from django.contrib.auth import login, logout, authenticate
from django.db import IntegrityError
from django.contrib.auth.decorators import login_required
from .forms import FormularioRegistro, PowerPointUploadForm
import os
import stripe
from django.conf import settings
from .models import Presentation, PerfilUsuario
from .forms import PowerPointUploadForm
from io import BytesIO
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from pptx import Presentation as PPTXPresentation
import os
from django.conf import settings
from django.http import FileResponse
from django.core.exceptions import ObjectDoesNotExist

stripe.api_key = settings.STRIPE_SECRET_KEY

# Create your views here.

def signup(request):
    if request.method == 'GET':
        return render(request, 'signup.html', {
            'form': FormularioRegistro()
        })
    else:
        form = FormularioRegistro(request.POST)
        if form.is_valid():
            username = form.cleaned_data['username']
            password = form.cleaned_data['password1']
            
            # Guardar temporalmente en sesión
            request.session['signup_username'] = username
            request.session['signup_password'] = password
            
            # Crear sesión de pago Stripe
            session = stripe.checkout.Session.create(
                payment_method_types=['card'],
                line_items=[{
                    'price_data': {
                        'currency': 'mxn',
                        'product_data': {
                            'name': 'Registro y acceso al sistema',
                        },
                        'unit_amount': 5000,  # $50.00 MXN
                    },
                    'quantity': 1,
                }],
                mode='payment',
                success_url=request.build_absolute_uri('/registro_exitoso/'),
                cancel_url=request.build_absolute_uri('/registro_cancelado/'),
                metadata={
                    'username': username
                }
            )
            return redirect(session.url, code=303)
        else:
            return render(request, 'signup.html', {
                'form': form,
                "error": "Formulario inválido"
            })


def signin(request):

    if request.method == 'GET':
        return render(request, 'signin.html', {
            'form': AuthenticationForm
        })
    else:
        user = authenticate(
            request, username=request.POST['username'], password=request.POST['password'])
        if user is None:
            return render(request, 'signin.html', {
                'form': AuthenticationForm,
                'error': 'Username or password is incorect'
            })
        else:
            login(request, user)
            return redirect('home')


@login_required
def signout(request):
    logout(request)
    return redirect('home')

def convert_pptx_to_pdf(pptx_path, pdf_path):
    ppt = PPTXPresentation(pptx_path)
    c = canvas.Canvas(pdf_path, pagesize=letter)
    width, height = letter

    for slide in ppt.slides:
        y_position = height - 50  # Margen superior

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines = shape.text.strip().split('\n')
                for line in lines:
                    if y_position < 50:
                        c.showPage()
                        y_position = height - 50
                    c.drawString(50, y_position, line)
                    y_position -= 20

        c.showPage()  # Nueva página para siguiente diapositiva

    c.save()

@login_required
def home(request):
    # Verifica si el usuario tiene un perfil y ha pagado
    if not hasattr(request.user, 'perfilusuario') or not request.user.perfilusuario.pagado:
        return redirect('crear_checkout')

    # Manejo de formulario para subir presentación
    if request.method == 'POST':
        form = PowerPointUploadForm(request.POST, request.FILES)
        if form.is_valid():
            presentation = form.save(commit=False)
            presentation.uploaded_by = request.user
            presentation.save()
            return redirect('home')
    else:
        form = PowerPointUploadForm()

    # Cargar presentaciones existentes (podrías filtrar por usuario si lo deseas)
    presentations = []
    for p in Presentation.objects.all().order_by('-uploaded_at'):
        if p.pptx_file and os.path.isfile(p.pptx_file.path):
            presentations.append(p)

    return render(request, 'home.html', {
        'form': form,
        'presentations': presentations
    })


@login_required
def delete_presentation(request, presentation_id):
    presentation = get_object_or_404(Presentation, pk=presentation_id)

    # Verificar que el usuario que intenta eliminar es el dueño
    if presentation.uploaded_by != request.user:
        return redirect('home')

    if request.method == 'POST':
        # Eliminar archivos físicos si existen
        if presentation.pptx_file and presentation.pptx_file.path:
            if os.path.exists(presentation.pptx_file.path):
                os.remove(presentation.pptx_file.path)

        # Eliminar el objeto de la base de datos
        presentation.delete()

        return redirect('home')


def root_redirect(request):
    if request.user.is_authenticated:
        return redirect('home')
    # Renderizamos el signin con el formulario para no usar redirección 302
    return render(request, 'signin.html', {
        'form': AuthenticationForm()
    })


@login_required
def crear_checkout(request):
    session = stripe.checkout.Session.create(
        payment_method_types=['card'],
        line_items=[{
            'price_data': {
                'currency': 'mxn',
                'product_data': {
                    'name': 'Acceso al sistema',
                },
                'unit_amount': 5000,  # $50.00 MXN
            },
            'quantity': 1,
        }],
        mode='payment',
        success_url='http://127.0.0.1:8000/pago_exitoso/',
        cancel_url='http://127.0.0.1:8000/pago_cancelado/',
        metadata={
            'user_id': request.user.id,
        }
    )
    return redirect(session.url, code=303)


@login_required
def pago_exitoso(request):
    perfil = request.user.perfilusuario
    perfil.pagado = True
    perfil.save()
    return redirect('home')

@login_required
def pago_cancelado(request):
    return render(request, 'pago_cancelado.html')

def inicio(request):
    return render(request, 'inicio.html')

def registro_exitoso(request):
    username = request.session.pop('signup_username', None)
    password = request.session.pop('signup_password', None)

    if username and password:
        try:
            user = User.objects.create_user(username=username, password=password)
            login(request, user)

            # Marca como pagado
            perfil = PerfilUsuario.objects.get(user=user)
            perfil.pagado = True
            perfil.save()

            return redirect('home')
        except IntegrityError:
            return redirect('/signup/?error=usuario_existente')
    else:
        return redirect('/signup/?error=datos_invalidos')

def registro_cancelado(request):
    # Limpiar sesión
    request.session.pop('signup_username', None)
    request.session.pop('signup_password', None)

    return render(request, 'signup.html', {
        'form': FormularioRegistro(),
        'error': 'El pago fue cancelado o fallido. Intenta de nuevo.'
    })
